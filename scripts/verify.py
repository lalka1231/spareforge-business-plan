"""Independent read-back checks; does not import the financial generator."""
from pathlib import Path
from decimal import Decimal as D
import csv,json,re,zipfile
from pptx import Presentation
ROOT=Path(__file__).resolve().parents[1]
def load(name):return list(csv.DictReader((ROOT/'data'/name).open(encoding='utf-8')))
def near(a,b):assert abs(D(str(a))-D(str(b)))<D('.02'),(a,b)
def main():
 rows=load('financial_model_5y.csv');summary=json.loads((ROOT/'data/finance_summary.json').read_text());assert len(rows)==6
 total=D(0)
 for i,r in enumerate(rows):
  n=lambda k:D(r[k])
  if i:
   near(n('revenue_rub'),n('treatment_ha_passes')*n('treatment_price_ex_vat_rub')+n('monitoring_ha_surveys')*n('monitoring_price_ex_vat_rub'))
   near(n('variable_cost_rub'),n('treatment_ha_passes')*n('treatment_variable_unit_rub')+n('monitoring_ha_surveys')*n('monitoring_variable_unit_rub'))
   near(n('ebitda_rub'),n('revenue_rub')-n('variable_cost_rub')-n('fixed_opex_rub'))
   near(n('usn_tax_rub'),n('revenue_rub')*D('.06'))
   near(n('accounting_profit_rub'),n('ebitda_rub')-n('depreciation_rub')-n('usn_tax_rub'))
   near(n('free_cash_flow_rub'),n('ebitda_rub')-n('usn_tax_rub')-n('additional_capex_rub')+n('reserve_return_rub'))
   near(n('vat_collected_rub'),n('vat_remitted_rub'))
   near(n('vat_collected_rub'),0 if i<4 else n('revenue_rub')*D('.05'))
  else:near(n('free_cash_flow_rub'),-D(9500000))
  total+=n('free_cash_flow_rub')/(D('1.24')**i)
 near(total,summary['base']['npv_rub'])
 near(sum(D(r['total_cost_rub']) for r in load('equipment.csv')),7000000)
 for r in load('unit_economics.csv'):
  near(D(r['price_ex_vat_rub'])-D(r['variable_cost_rub']),r['contribution_before_usn_rub'])
  near(D(r['price_ex_vat_rub'])*D('.94')-D(r['variable_cost_rub']),r['contribution_after_usn_rub'])
 monthly=load('monthly_cashflow_2027.csv');assert len(monthly)==12
 near(sum(D(r['net_operating_cf_rub']) for r in monthly),rows[1]['free_cash_flow_rub'])
 near(min(D(r['closing_cash_rub']) for r in monthly),950000)
 assert [r['month'] for r in monthly if D(r['share'])>0]==[f'2027-{i:02}' for i in range(4,10)]
 for case in summary['yearly_scenarios'].values():
  independently=sum(D(str(r['free_cash_flow_rub']))/D('1.24')**i for i,r in enumerate(case))
  key=next(k for k,v in summary['yearly_scenarios'].items() if v is case)
  near(independently,next(r['npv_rub'] for r in summary['sensitivity'] if r['scenario']==key))
 plan=load('project_plan.csv');by={r['id']:r for r in plan};assert len(by)==13
 for r in plan:
  assert len(r['A'].split(';'))==1
  assert int(r['ef'])-int(r['es'])==int(r['duration_weeks'])
  for p in r['predecessors'].split(';'):
   if p:assert int(r['es'])>=int(by[p]['ef'])
 assert by['K']['start_date']=='2027-04-12' and by['K']['finish_date']=='2027-09-12'
 doc=(ROOT/'docs/FINAL_BUSINESS_PLAN.md').read_text();ids=[int(n) for n in re.findall(r'^## (\d+)\.',doc,re.M)];assert ids==list(range(1,23))
 src=json.loads((ROOT/'data/source_map.json').read_text());valid={r['id'] for r in src};cited={int(x) for x in re.findall(r'\[(\d+)\]',doc)};assert cited<=valid
 assert all((ROOT/f'docs/evidence/{i:02d}.txt').exists() for i in cited)
 ppt=ROOT/'presentations/AGROVECTOR_22_SLIDES.pptx';prs=Presentation(ppt);assert len(prs.slides)==22
 for s in prs.slides:
  assert s.notes_slide.notes_text_frame.text.strip()
  for sh in s.shapes:
   assert sh.left>=0 and sh.top>=0 and sh.left+sh.width<=prs.slide_width+20 and sh.top+sh.height<=prs.slide_height+20
 texts=[]
 with zipfile.ZipFile(ppt) as z:
  texts=[z.read(n).decode('utf-8') for n in z.namelist() if n.endswith('.xml')]
 assert 'spareforge' not in '\n'.join(texts).lower()
 for folder in ['docs','data']:
  for path in (ROOT/folder).rglob('*'):
   if path.is_file() and path.suffix in ['.md','.csv','.json','.txt']:
    if path.name=='FINAL_REVIEW.md':continue
    text=path.read_text();assert 'spareforge' not in text.lower(),path
 assert not re.search(r'учебн[а-я]* (?:бизнес|проект)|студент|преподавател|задание по дисциплине',doc,re.I)
 outline=json.loads((ROOT/'presentations/slide_outline.json').read_text());assert not outline['shape_bounds_errors'];assert not outline['body_font_errors']
 result={'verified':True,'sections':22,'slides':22,'project_tasks':13,'sources_registered':len(valid),'sources_cited':len(cited),'npv_independent_rub':float(total),'capex_rub':7000000,'initial_funding_rub':9500000,'monthly_min_cash_rub':950000,'csv_rows':{x:len(load(x)) for x in ['financial_model_5y.csv','assumptions.csv','competitors.csv','project_plan.csv','sensitivity.csv']},'visual_rendering':False,'limitations':['PowerPoint rendering unavailable: LibreOffice and poppler missing','Source verification is evidence-based retrieval, not legal opinion','TAM/SAM and future shares are scenarios, not independently measured markets']}
 (ROOT/'data/validation.json').write_text(json.dumps(result,ensure_ascii=False,indent=2),encoding='utf-8');print(json.dumps(result,ensure_ascii=False,indent=2))
if __name__=='__main__':main()
