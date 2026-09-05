"""Build the editable 22-slide AgroVector investor deck; structural QA only.
Run: /tmp/agrovector-venv/bin/python scripts/build_presentation.py
Notes refresh from the numbered final AgroVector document on every run.
"""
from pathlib import Path
import ast
import csv
import json
import re
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.xmlchemy import OxmlElement

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / 'presentations'
FOREST, DARK, LIME, WHITE, MIST, MUTED, RED = '164C38', '092A20', 'C8F26B', 'FFFFFF', 'EFF4EF', 'AFC5B8', 'C35442'
W, H = 13.333333, 7.5
FONT = 'Arial'

def rgb(v): return RGBColor.from_string(v)
def inch(v): return Inches(v)
def million(v, digits=2): return f'{v / 1e6:.{digits}f}'.replace('.', ',')
def load_csv(name):
    with (ROOT / 'data' / name).open(encoding='utf-8-sig', newline='') as f:
        return list(csv.DictReader(f))

def get_titles():
    tree = ast.parse((ROOT / 'scripts/build_documents.py').read_text(encoding='utf-8'))
    for node in tree.body:
        if isinstance(node, ast.Assign) and any(isinstance(t, ast.Name) and t.id == 'TITLES' for t in node.targets):
            result = ast.literal_eval(node.value)
            assert len(result) == 22
            return result
    raise ValueError('TITLES missing')

def sections_for_notes():
    result = {}
    for path in [Path('/tmp/agro_sections.json'), Path('/tmp/agro_parent_sections.json'), ROOT / 'data/sections.json']:
        if path.exists():
            d = json.loads(path.read_text(encoding='utf-8'))
            # Repo sections are only safe once the final AgroVector document exists.
            if path.parent == ROOT / 'data':
                continue
            result.update(d)
    final = ROOT / 'docs/FINAL_BUSINESS_PLAN.md'
    if final.exists():
        text = final.read_text(encoding='utf-8')
        if re.match(r'#\s+АгроВектор\b', text):
            matches = list(re.finditer(r'^##\s+(\d+)\.\s+(.+)$', text, re.M))
            for i, m in enumerate(matches):
                result[m[1]] = text[m.end():matches[i+1].start() if i+1 < len(matches) else len(text)].strip()
    return result

class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = inch(W), inch(H)
        self.titles = get_titles()
        self.sections = sections_for_notes()
        self.sources = {}
        self.note_extras = {}

    def box(self, s, x, y, w, h, fill=WHITE, radius=False):
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, inch(x), inch(y), inch(w), inch(h))
        shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill); shape.line.fill.background()
        if radius:
            shape.adjustments[0] = 0.08
        return shape

    def text(self, s, text, x, y, w, h, size=18, color=DARK, bold=False):
        shape = s.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
        tf = shape.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = inch(.015)
        tf.margin_top = tf.margin_bottom = 0
        for i, line in enumerate(str(text).split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line; p.font.name = FONT; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = rgb(color)
            p.space_after = Pt(5); p.space_before = Pt(0)
        return shape

    def slide(self, n, label='[ДОПУЩЕНИЕ]', sources=()):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb(DARK)
        self.box(s, .45, .43, .1, .49, LIME)
        self.text(s, self.titles[n-1], .72, .39, 11.85, .97, 26 if len(self.titles[n-1]) > 48 else 30, WHITE, True)
        self.text(s, label, .72, 1.27, 11.9, .36, 15, LIME)
        self.text(s, 'АГРОВЕКТОР / инвестиционная гипотеза / 2027–2031', .55, 7.12, 11, .23, 10, MUTED)
        self.text(s, f'{n:02d} / 22', 12.05, 7.09, .8, .28, 11, LIME, True)
        if sources:
            self.sources[n] = list(sources)
            self.text(s, 'Источники: ' + ' · '.join(sources), .55, 6.64, 12.2, .44, 9, MUTED)
        return s

    def card(self, s, x, y, w, h, title, body, number=None, size=18):
        self.box(s, x, y, w, h, WHITE, True)
        self.box(s, x+.2, y+.2, .38, .06, LIME)
        offset = .0
        if number is not None:
            self.text(s, number, x+.22, y+.39, w-.44, .72, 35, FOREST, True)
            offset = .77
        self.text(s, title, x+.22, y+.42+offset, w-.44, .64, 20, FOREST, True)
        self.text(s, body, x+.22, y+1.08+offset, w-.44, h-1.16-offset, size)

    def banner(self, s, text, y=6.03, color=LIME):
        self.box(s, .55, y, 12.23, .49, color, True)
        self.text(s, text, .72, y+.09, 11.9, .35, 16, DARK, True)

    def table(self, s, headers, rows, x=.55, y=1.9, w=12.23, h=3.8, widths=None, size=16):
        shape = s.shapes.add_table(len(rows)+1, len(headers), inch(x), inch(y), inch(w), inch(h))
        t = shape.table
        if widths:
            for c, width in zip(t.columns, widths): c.width = inch(width)
        for ri, row in enumerate([headers]+rows):
            for ci, value in enumerate(row):
                c = t.cell(ri, ci); c.text = str(value)
                c.fill.solid(); c.fill.fore_color.rgb = rgb(FOREST if ri == 0 else WHITE if ri%2 else MIST)
                c.margin_left=c.margin_right=inch(.1); c.margin_top=c.margin_bottom=inch(.05)
                for p in c.text_frame.paragraphs:
                    p.font.name=FONT; p.font.size=Pt(size); p.font.color.rgb=rgb(WHITE if ri==0 else DARK); p.font.bold=(ri==0)
                    p.space_after=Pt(0)
        return shape

    def chart(self, s, categories, series, x, y, w, h, kind=XL_CHART_TYPE.COLUMN_CLUSTERED, colors=None):
        self.box(s, x, y, w, h, WHITE, True)
        data=CategoryChartData(); data.categories=categories
        for name, values in series: data.add_series(name, values)
        chart=s.shapes.add_chart(kind, inch(x+.15), inch(y+.12), inch(w-.3), inch(h-.24), data).chart
        chart.font.name=FONT; chart.font.size=Pt(15)
        chart.has_legend=len(series)>1
        if chart.has_legend:
            chart.legend.position=XL_LEGEND_POSITION.BOTTOM; chart.legend.include_in_layout=False
            chart.legend.font.size=Pt(15)
        chart.chart_style=10
        palette=colors or [FOREST, '7EB94A', '739DAB', 'BC9255']
        for i, ser in enumerate(chart.series):
            ser.format.fill.solid(); ser.format.fill.fore_color.rgb=rgb(palette[i%len(palette)])
            ser.format.line.color.rgb=rgb(palette[i%len(palette)])
            ser.format.line.width=Pt(2)
        chart.category_axis.tick_labels.font.size=Pt(15)
        chart.value_axis.tick_labels.font.size=Pt(15)
        chart.value_axis.has_major_gridlines=True
        return chart

    def finish(self):
        assert len(self.prs.slides)==22
        for n,s in enumerate(self.prs.slides,1):
            texts=[]
            for sh in s.shapes:
                if sh.has_text_frame: texts.append(sh.text)
                if sh.has_table: texts.extend(' | '.join(c.text for c in row.cells) for row in sh.table.rows)
            notes=self.sections.get(str(n), '\n'.join(texts))
            notes+='\n\nКомментарий к слайду:\n'+self.note_extras.get(n,'Редактируемые объекты. Финансовые показатели — расчет модели, входные данные — допущения; проект не является действующим подтвержденным бизнесом.')
            if n in self.sources: notes+='\n\nВнешние URL:\n'+'\n'.join(self.sources[n])
            s.notes_slide.notes_text_frame.text=notes
        OUT.mkdir(exist_ok=True)
        path=OUT/'AGROVECTOR_22_SLIDES.pptx'; self.prs.save(path)
        # Independently reopen the saved package and extract its actual content.
        reopened=Presentation(path); allslides=[]; bounds=[]; fontsize=[]; charts=0; tables=0
        for n,s in enumerate(reopened.slides,1):
            entry={'slide':n,'title':self.titles[n-1],'texts':[],'tables':[],'charts':[],'speaker_notes':s.notes_slide.notes_text_frame.text,'sources':self.sources.get(n,[])}
            for sh in s.shapes:
                if sh.left<0 or sh.top<0 or sh.left+sh.width>reopened.slide_width+2 or sh.top+sh.height>reopened.slide_height+2: bounds.append([n,sh.name])
                if sh.has_text_frame:
                    entry['texts'].append(sh.text)
                    for p in sh.text_frame.paragraphs:
                        if p.font.size and p.font.size.pt < (11 if n==16 else 15) and sh.top < inch(6.6): fontsize.append([n,sh.text,p.font.size.pt])
                if sh.has_table:
                    tables+=1; entry['tables'].append([[c.text for c in r.cells] for r in sh.table.rows])
                    for row in sh.table.rows:
                        for cell in row.cells:
                            for p in cell.text_frame.paragraphs:
                                if p.font.size.pt < (11 if n==16 else 15): fontsize.append([n,cell.text,p.font.size.pt])
                if sh.has_chart:
                    charts+=1; chart=sh.chart
                    entry['charts'].append({'categories':[c.label for c in chart.plots[0].categories],'series':[{'name':ser.name,'values':list(ser.values)} for ser in chart.series]})
            allslides.append(entry)
        assert len(reopened.slides)==22 and not bounds and not fontsize, (bounds,fontsize)
        assert all(s['speaker_notes'].strip() for s in allslides)
        report={'slide_count':len(allslides),'aspect_ratio':'16:9','shape_bounds_errors':bounds,'body_font_errors':fontsize,'editable_charts':charts,'editable_tables':tables,'speaker_notes_count':len(allslides),'verification':'Structural read-back only; no visual rendering performed.','soffice_available':bool(shutil.which('soffice')),'slides':allslides}
        (OUT/'slide_outline.json').write_text(json.dumps(report,ensure_ascii=False,indent=2),encoding='utf-8')
        print(json.dumps({k:v for k,v in report.items() if k!='slides'},ensure_ascii=False,indent=2))
        print(path)


def build():
    d=Deck(); f=json.loads((ROOT/'data/finance_summary.json').read_text()); base=f['base']; years=[r for r in f['yearly_scenarios']['base'] if r['period']>0]
    tasks=load_csv('project_plan.csv'); assert len(tasks)==13
    byid={t['id']:t for t in tasks}; short={'A':'Правовой gate','B':'Спрос / заявки','C':'ООО и договоры','D':'Финансирование','E':'Закупка','F':'Обучение','G':'Готовность БВС','H':'SOP / CRM','I':'Испытания','J':'Договоры сезона','M':'Ожидание окна','K':'Сезон','L':'Аудит'}
    critical=[t['id'] for t in tasks if t['critical'].lower()=='true']; duration=max(int(t['ef']) for t in tasks)
    s=d.slide(1,'[ДОПУЩЕНИЕ] Концепция сервиса  •  [РАСЧЕТ] Финансы базового сценария')
    d.text(s,'АгроВектор',.65,1.93,7,1,48,WHITE,True)
    d.text(s,'Мониторинг посевов и точечное внесение\nс проверяемым заданием и отчетом',.7,3.0,6.5,1.15,25,WHITE)
    d.text(s,'Одна бригада • Татарстан — условная зона\nГражданская переподготовка специалистов БПЛА',.7,4.53,6.2,1.0,18,MUTED)
    d.card(s,7.25,1.92,2.65,3.73,'млн ₽','Стартовая потребность',million(base['initial_investment_rub'],1))
    d.card(s,10.12,1.92,2.65,3.73,'млн ₽ NPV','Ставка 24%; горизонт 5 лет',million(base['npv_rub']))
    d.banner(s,'Решение: поэтапное финансирование после правовой проверки и подтверждения спроса')

    s=d.slide(2)
    d.table(s,['SMART-результат','Срок','Подтверждение'],[
        ['Правовая применимость','До go/no-go 2026','Письменная карта ограничений'],
        ['20 интервью; заявки на 6 000 га-проходов','До финансирования','Реестр клиентов и полей'],
        ['10 сезонных клиентов; 2 оператора','До 12.04.2027','Договоры; оценка компетенций'],
        ['6 000 га-проходов + 4 000 га-обследований','Сезон 2027','Акты и журналы; пилот внутри объема'],
        ['Добровольная гражданская переподготовка','До самостоятельных работ','Наставник; единые критерии допуска']
    ],h=3.85,widths=[5.15,2.8,4.28],size=17)
    d.banner(s,'Коммерческие цели не отменяют право остановить небезопасное задание')

    s=d.slide(3,'[ДОПУЩЕНИЕ] Продукт и стратегическая канва: экспертная шкала 1–5, не исследование рынка')
    d.text(s,'Обследование → карта → внесение → отчет',.65,1.82,12,.5,23,WHITE,True)
    d.chart(s,['Локальный\nобъем','Геометрия','Связка\nс картой','Отчет','Без CAPEX\nклиента','Доступность','Широта\nтехнологий'],[
        ('АгроВектор',[5,5,5,5,5,3,3]),('Наземный подрядчик',[2,2,2,3,5,3,4]),('Своя техника',[3,3,3,3,1,5,4]),('Узкий БПЛА-подрядчик',[4,4,3,3,5,3,3])],.55,2.46,12.23,3.48,XL_CHART_TYPE.LINE)
    d.banner(s,'Препараты клиента; без обещания роста урожая. Орошение и продажа дронов не входят')

    s=d.slide(4)
    blocks=[('1 / Сегменты','КФХ и средние хозяйства; локальные сложные участки'),('2 / Ценность','Обследование, карта задания, внесение и отчет'),('3 / Каналы','Директор; агрономы; демонстрации; рекомендации'),('4 / Отношения','Один контакт; календарь сезона; разбор отклонений'),('5 / Доходы','₽/га-проход и ₽/га-обследование; без двойного счета'),('6 / Ресурсы','Бригада; основной и резервный БВС; мониторинг'),('7 / Действия','Правовая проверка; выполнение; контроль; сервис'),('8 / Партнеры','Агроном; обучение; сервис; юрист; бухгалтер'),('9 / Затраты','Техника; труд; эксплуатация; услуги; ликвидность')]
    for i,(a,b) in enumerate(blocks):
        x=.55+(i%3)*4.15; y=1.85+(i//3)*1.51
        d.box(s,x,y,3.93,1.35,WHITE,True); d.text(s,a,x+.18,y+.14,3.57,.35,19,FOREST,True); d.text(s,b,x+.18,y+.58,3.57,.68,16)

    s=d.slide(5,'[РАСЧЕТ] на [ДОПУЩЕНИЕ]: TAM / SAM / SOM — сценарные, не внешняя статистика',sources=['https://www.vesti.ru/article/4893778','https://www.rshb.ru/news/16072026-000004'])
    for i,(num,title,body) in enumerate([('1,245 млрд','TAM / ₽ в год','2 000 хозяйств × 2 000 га\n15% пригодной площади'),('124,5 млн','SAM / ₽ в год','200 хозяйств × 2 000 га\nУсловный доступный кластер'),('8,8 млн','SOM / первый сезон','6 000 га-проходов × 1 300 ₽\n+ 4 000 обследований × 250 ₽')]):
        d.card(s,.55+4.15*i,1.9,3.93,2.78,title,body,num,16)
    d.text(s,'TAM = 2 000 × 2 000 × 15% × (1,5 × 1 300 + 0,5 × 250); SAM — для 200 хозяйств.',.65,4.9,12,.65,17,WHITE)
    d.text(s,'Воронка: 100 контактов → 40 разговоров → 20 предложений → 10 клиентов.\n[ФАКТ] Публикации об отрасли имеют разные оценки; не используются как база TAM.',.65,5.62,12,.84,16,MUTED)
    d.note_extras[5]='TAM/SAM — собственная гипотеза снизу вверх. Вести: около 600 тыс. га обработано в 2025, со слов министра; РСХБ: 781 тыс. га в 2025 и прогноз 1058 тыс. га в 2026. Различия определений не разрешены; ряды не смешиваются. Все числа рынка на слайде сценарные. Бюджет маркетинга 240 тыс. рублей; денежный CAC 24 тыс. рублей без зарплаты директора. URL — контекст отрасли, не источник TAM/SAM.'

    s=d.slide(6,'[ФАКТ] Публичные предложения  •  [ДОПУЩЕНИЕ] Тариф АгроВектора',sources=['https://agrdron.ru/','https://beeuav.ru/','https://xag-ap.ru/spraying.php','https://germes-technology.ru/'])
    d.table(s,['Исполнитель','Публичное предложение','Цена, ₽/га','Доля рынка'],[
        ['Агродроны ЮФО','Обработка полей; ЮФО','900 — средняя заявленная','н/д'],
        ['BeeUAV','Аренда и обработка','н/д; по запросу','н/д'],
        ['Агропрофи / XAG','Обработка и посев; Поволжье','н/д','н/д'],
        ['Гермес','Краснодарский край, Адыгея','н/д','н/д'],
        ['Аграс / аграс.рф','Аренда техники / пилота','н/д','н/д'],
        ['АгроВектор [ДОПУЩЕНИЕ]','Карта + выполнение + отчет','1 300 без НДС','н/д']
    ],widths=[3.25,4.05,3.43,1.5],h=4.1,size=15)
    d.banner(s,'900 ₽ — не нижняя граница и не сопоставимое КП. Нужен единый состав работ для сравнения')
    d.note_extras[6]='Агродроны ЮФО: «Стоимость зависит от вида обработки и объёма участка. В среднем 900 рублей за гектар. Итоговая цена рассчитывается индивидуально». Состав, НДС, логистика и препараты требуют запроса КП. Доли всех участников неизвестны. BeeUAV рекламирует Татарстан; адрес на сайте — Московская область, локальное присутствие требует проверки. Канва слайда 3 не является оценкой перечисленных компаний.'

    s=d.slide(7)
    for i,(title,body) in enumerate([('S / Конструкция сервиса','Мониторинг + задание + отчет\nРезервный аппарат\n→ Единая приемка и готовность резерва'),('W / Ограничения','Нет подтвержденного портфеля\nОдна бригада; экономный ФОТ\n→ Заявки до закупки, лимит нагрузки'),('O / Возможности','Локальные задачи хозяйств\nГражданская переподготовка\n→ Платный пилот и наставничество'),('T / Угрозы','Запреты; погода; ценовое давление\nПросрочка платежей\n→ No-go, резерв времени и аванс')]):
        d.card(s,.55+(i%2)*6.24,1.88+(i//2)*2.23,5.99,2.06,title,body,size=17)
    d.banner(s,'Стратегия WT: не расширять парк до подтверждения правовой доступности, кадров и экономики',6.43)

    s=d.slide(8)
    d.table(s,['Сила','Уровень гипотезы','Ответ проекта / что проверить'],[
        ['Соперничество','Высокий','Приемка и сроки; сопоставимые КП'],['Новые участники','Средний','Повторные договоры; стоимость готовности'],['Заменители','Высокий','Отбор задач; полная стоимость альтернатив'],['Покупатели','Высокий','Лимит концентрации; причины отказов'],['Поставщики','Средний–высокий','Запчасти, резерв, сервис; реальные сроки ремонта']
    ],widths=[3.4,2.55,6.28],h=3.85,size=18)
    d.banner(s,'Устойчивое конкурентное преимущество до полевой проверки не подтверждено')

    legal_urls=['https://base.garant.ru/414829385/','http://publication.pravo.gov.ru/document/0001202608280041']
    legal_section=d.sections.get('9','')
    discovered=re.findall(r'https?://[^\s)>\]]+',legal_section)
    if discovered: legal_urls=list(dict.fromkeys(discovered))[:3]
    s=d.slide(9,'[ДОПУЩЕНИЕ] ООО и порядок проверки; не юридическое заключение',sources=legal_urls)
    d.card(s,.55,1.88,3.9,3.83,'ООО','Договорный контур и долевое финансирование.\n\nСтатус БВС определяется по максимальной взлетной массе, не по массе пустого аппарата.',size=18)
    d.table(s,['Отдельный gate','Что подтвердить до работы'],[['БВС','Учет / регистрация и применимые документы'],['Полет','Территория, маршрут и условия использования'],['АХР','Правомерность конкретной схемы выполнения'],['Препарат','Допустимость технологии и регламент применения'],['Регион','Действующий запрет / исключения, если применимы']],x=4.7,w=8.08,h=3.83,widths=[1.9,6.18],size=17)
    d.banner(s,'[ФАКТ] С 01.09.2026 — ФАП-303. Санитарные условия требуют отдельной сверки; при сомнении NO-GO')
    d.note_extras[9]='До закупки требуется письменная карта применимых норм от юриста, а перед выездом — повторная проверка. Базовые URL Росавиации и портала правовой информации — точки входа для проверки, не подтверждение выданного разрешения. При наличии окончательной секции 9 ее конкретные внешние ссылки автоматически подхватываются. Не смешивать массу пустого БВС, полезную нагрузку и максимальную взлетную массу. Мониторинг-only не доказан как прибыльная замена базовой модели.'

    s=d.slide(10)
    for i,(a,b) in enumerate([('Высокая власть / высокий интерес','Учредитель, инвестор, клиенты и агрономы\n→ План-факт; договоры; приемка\nВладелец: директор'),('Высокая власть / переменный интерес','Уполномоченные органы\n→ Документы по конкретной территории\nВладелец: внешний юрист'),('Средняя власть / высокий интерес','Операторы, учебный партнер, агроном\n→ Инструктаж, нагрузка, компетенции\nВладелец: операционный руководитель'),('Низкая власть / высокий или переменный','Кандидаты после СВО, соседи полей\n→ Добровольность; канал обращений\nВладелец: директор')]):
        d.card(s,.55+(i%2)*6.24,1.88+(i//2)*2.23,5.99,2.06,a,b,size=16)
    d.text(s,'Сервис техники: средняя власть / средний интерес; регламент реакции — операционный руководитель.',.65,6.45,12,.55,15,MUTED)

    s=d.slide(11,'[ДОПУЩЕНИЕ] Производительность  •  [РАСЧЕТ] Календарная емкость бригады')
    d.card(s,.55,1.88,3.92,3.63,'га-проходов / потолок','100 эффективных дней × 130 га\nОсновной + резервный аппарат\nРезерв не удваивает мощность','13 000',17)
    d.card(s,4.7,1.88,3.92,3.63,'га-обследований','20 отдельных дней × 500 га\nБез одновременных полетов\nОдной бригадой — по очереди','10 000',17)
    d.card(s,8.85,1.88,3.92,3.63,'дня резерва','154 календарных дня сезона\n− 100 дней внесения\n− 20 дней мониторинга','34',17)
    d.text(s,'Проверка → обследование → подготовка → выполнение → контроль и акт',.65,5.66,12,.5,20,WHITE,True)
    d.banner(s,'Это напряженный сценарный потолок. Хронометраж и достаточность резерва не подтверждены',6.27)

    s=d.slide(12)
    d.text(s,'3 FTE',.65,1.92,4,1,44,LIME,True)
    d.text(s,'3,0 млн ₽ / год — бюджет штатного труда; требует проверки рынка и начислений',3.5,2.02,9.1,.86,20,WHITE)
    d.table(s,['Штатная роль','Функции'],[['Директор','Продажи, договоры, деньги и партнеры'],['Оператор / операционный руководитель','Календарь, техника, безопасность и наставник'],['Второй оператор / помощник','Подготовка; мониторинг в отдельных окнах']],y=3.02,h=1.91,widths=[4.3,7.93],size=18)
    d.text(s,'Вне штата: агроном, юрист, бухгалтер и учебный партнер — оплачиваемые услуги.',.65,5.15,12,.6,18,MUTED)
    d.banner(s,'После СВО: добровольный отбор → гражданское обучение → наставник → проверка компетенций')
    d.text(s,'Одинаковые требования для всех; опыт не дает автоматического допуска. Личные данные — конфиденциально.',.65,6.65,12,.39,15,MUTED)

    s=d.slide(13,'[ДОПУЩЕНИЕ] WBS из data/project_plan.csv; каждый пакет имеет результат')
    groups=[('1 / Проверка модели','1.1 A — правовая карта\n1.2 B — интервью и заявки\n1.3 C — ООО и договоры'),('2 / Ресурсы','2.1 D — инвестиционное решение\n2.2 E — приемка техники\n2.3 F — обучение\n2.4 G — готовность БВС'),('3 / Предсезонная готовность','3.1 H — SOP, CRM, карты\n3.2 I — приемочные испытания\n3.3 J — договоры сезона\n3.4 M — ожидание окна'),('4 / Исполнение и контроль','4.1 K — коммерческий сезон\n4.2 L — аудит и решение')]
    for i,(a,b) in enumerate(groups): d.card(s,.55+(i%2)*6.24,1.88+(i//2)*2.23,5.99,2.06,a,b,size=17)
    d.text(s,'Закупка ≠ разрешение на работы. Полевой пилот 100 га включен в апрельский сезон, не прибавлен сверху.',.65,6.48,12,.52,15,MUTED)

    s=d.slide(14,'[РАСЧЕТ] Гант из CSV  •  [ДОПУЩЕНИЕ] Даты и длительности — план, не выданные разрешения')
    x0,y0,plotw=4.13,2.15,8.15
    d.box(s,.55,1.86,12.23,4.55,WHITE,True)
    for week in range(0,duration+1,5):
        x=x0+plotw*week/duration
        d.box(s,x,y0,.008,3.97,MIST); d.text(s,str(week),x-.15,1.91,.4,.25,15,FOREST)
    rowh=.29
    for i,t in enumerate(tasks):
        y=y0+i*rowh
        d.text(s,t['id']+'  '+short[t['id']],.7,y,3.3,.29,15)
        d.box(s,x0+plotw*int(t['es'])/duration,y+.045,max(.035,plotw*int(t['duration_weeks'])/duration),.18,FOREST if t['critical'].lower()=='true' else '91B663',True)
    d.text(s,f'Недели от {tasks[0]["start_date"]}  •  Конец {byid["L"]["finish_date"]}  •  Темный = критический путь',.7,6.04,11.7,.33,15,FOREST)
    d.text(s,f'Сезон K: {byid["K"]["start_date"]} — {byid["K"]["finish_date"]}; пакет M: {byid["M"]["duration_weeks"]} недель ожидания окна.',.65,6.59,12,.45,16,MUTED)

    s=d.slide(15,'[РАСЧЕТ] Сеть finish-to-start из CSV; длительности — календарные недели')
    # Directed acyclic graph laid out by dependency depth, not a decorative chain.
    depth={}
    for t in tasks:
        ps=[p for p in t['predecessors'].split(';') if p]
        depth[t['id']]=max((depth[p]+1 for p in ps),default=0)
    cols={}
    for tid,dep in depth.items(): cols.setdefault(dep,[]).append(tid)
    positions={}
    step=11.9/(max(depth.values())+1)
    for dep,tids in cols.items():
        for j,tid in enumerate(tids): positions[tid]=(.65+dep*step,2.55+j*1.75)
    nodew=.96; nodeh=.96
    for t in tasks:
        tx,ty=positions[t['id']]
        for p in filter(None,t['predecessors'].split(';')):
            px,py=positions[p]
            conn=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,inch(px+nodew),inch(py+nodeh/2),inch(tx),inch(ty+nodeh/2))
            conn.line.color.rgb=rgb(LIME if p in critical and t['id'] in critical else MUTED); conn.line.width=Pt(1.7)
            arrow=OxmlElement('a:tailEnd'); arrow.set('type','triangle'); conn.line._get_or_add_ln().append(arrow)
    for t in tasks:
        x,y=positions[t['id']]; d.box(s,x,y,nodew,nodeh,LIME if t['id'] in critical else WHITE,True)
        d.text(s,t['id'],x+.09,y+.11,.78,.36,24,FOREST,True)
        d.text(s,t['duration_weeks']+' нед.',x+.07,y+.59,.84,.3,15)
    d.text(s,'Критический путь: '+' → '.join(critical),.65,5.62,12,.44,20,WHITE,True)
    d.banner(s,f'{duration} недель от проверки до аудита; M — явная работа ожидания, не свободный резерв',6.27)
    d.note_extras[15]='Связи и длительности полностью извлечены из project_plan.csv.\n'+'\n'.join(f"{t['id']} {t['task']}: предшественники {t['predecessors'] or 'нет'}, ES {t['es']}, EF {t['ef']}, LS {t['ls']}, LF {t['lf']}, резерв {t['slack_weeks']}" for t in tasks)

    s=d.slide(16,'[ДОПУЩЕНИЕ] R — выполняет; A — принимает; C — консультирует; I — информируется')
    abbrev={'Операционный руководитель':'Опер. рук.','Учебный партнер':'Учеб. партнер','Бухгалтер':'Бухгалтер'}
    def role(v):
        for a,b in abbrev.items(): v=v.replace(a,b)
        return v.replace(';',', ')
    rows=[[t['id']+' '+short[t['id']],role(t['R']),role(t['A']),role(t['C']),role(t['I'])] for t in tasks]
    d.table(s,['Работа','R','A','C','I'],rows,y=1.85,h=4.65,widths=[2.62,2.0,1.9,3.73,1.98],size=12)
    d.text(s,'Все 13 задач; один A в каждой строке. Продажи = директор; опер. рук. = один из двух операторов.',.65,6.63,12,.42,15,MUTED)
    d.note_extras[16]='Матрица извлечена из CSV без изменения распределения ответственности. Опер. рук. — операционный руководитель; учеб. партнер — учебный партнер.\n'+'\n'.join(f"{t['id']} {t['task']}: R {t['R']}; A {t['A']}; C {t['C']}; I {t['I']}" for t in tasks)

    s=d.slide(17,'[РАСЧЕТ] Базовый сценарий; млн ₽, выручка без НДС. Входные данные — допущения')
    metrics=[('Выручка','revenue_rub'),('EBITDA','ebitda_rub'),('Свободный CF','free_cash_flow_rub')]
    d.chart(s,[str(r['year']) for r in years],[(name,[r[key]/1e6 for r in years]) for name,key in metrics],.55,1.86,7.15,4.11)
    d.table(s,['Год','Выручка','EBITDA','CF'],[[r['year']]+[million(r[k]) for _,k in metrics] for r in years],x=7.93,y=1.86,w=4.85,h=3.06,widths=[.95,1.3,1.3,1.3],size=15)
    d.text(s,'CF₀ = −9,50 млн ₽\nCF 2031 включает возврат\nрезерва 2,00 млн ₽.',8.03,5.1,4.55,.93,17,WHITE)
    d.banner(s,'Продажа оборудования в терминале = 0. Балансовая стоимость не превращена в денежный поток',6.27)

    s=d.slide(18,'[РАСЧЕТ] NPV на 5 лет; ставка дисконтирования 24%; млн ₽')
    labels={'base':'База','price_minus_10pct':'Цена −10%','volume_minus_25pct':'Объем −25%','variable_plus_20pct':'Переменные +20%','one_year_delay':'Задержка на год','treatment_price_900':'Внесение 900 ₽/га'}
    scenarios=f['sensitivity']
    d.chart(s,[labels[t['scenario']] for t in scenarios],[('NPV',[t['npv_rub']/1e6 for t in scenarios])],.55,1.87,7.03,4.13,XL_CHART_TYPE.BAR_CLUSTERED)
    d.table(s,['Сценарий','NPV, млн ₽'],[[labels[t['scenario']],million(t['npv_rub'],3)] for t in scenarios],x=7.81,y=1.87,w=4.97,h=2.92,widths=[3.42,1.55],size=16)
    d.text(s,f'Окупаемость базы: {base["simple_payback_years"]:.2f} года\nДисконтированная: {base["discounted_payback_years"]:.2f} года'.replace('.',','),7.93,5.05,4.62,.95,18,WHITE)
    d.banner(s,'Снижение объема и задержка запуска дают отрицательный NPV. Цена −10% почти обнуляет запас',6.27)

    s=d.slide(19,'[ДОПУЩЕНИЕ] Бюджет до коммерческих предложений  •  [РАСЧЕТ] Сумма потребности')
    equipment=load_csv('equipment.csv')
    eqnames=['2 агродрона с АКБ / зарядкой','Транспорт / прицеп','Мониторинговый БВС','Подготовка состава / безопасность','ИТ','Капитальные запчасти']
    d.table(s,['CAPEX','млн ₽'],[[name,million(float(t['total_cost_rub']),1)] for name,t in zip(eqnames,equipment)],x=.55,y=1.9,w=7.2,h=3.92,widths=[5.5,1.7],size=17)
    for i,(num,title) in enumerate([('7,0','CAPEX'),('0,5','Стартовые расходы'),('2,0','Оборотный резерв')]):
        y=1.9+i*1.33; d.box(s,8.0,y,4.78,1.13,WHITE,True); d.text(s,num,8.2,y+.21,1.25,.64,32,FOREST,True); d.text(s,title,9.62,y+.29,2.95,.64,19,FOREST)
    d.banner(s,'Итого 9,5 млн ₽. Резерв — ликвидность, не повторный расход и не бюджет расширения',6.27)

    s=d.slide(20)
    d.card(s,.55,1.9,5.99,3.57,'млн ₽ / учредитель','Подтвердить наличие денег\nи оформить вклад','2,5',20)
    d.card(s,6.79,1.9,5.99,3.57,'млн ₽ / внешний инвестор','Долевое финансирование\nпосле проверки условий','7,0',20)
    d.text(s,'Долг и гранты не включены. Деньги пока не подтверждены.\nОценка бизнеса, корпоративные права, выход и доля инвестора — предмет соглашения.',.65,5.66,12,.83,18,WHITE)
    d.banner(s,'Это суммы источников, НЕ доли в капитале. Крупные закупки — только после инвестиционного gate',6.61)

    s=d.slide(21,'[РАСЧЕТ] Финансовые KPI базы  •  [ДОПУЩЕНИЕ] Плановые операционные пороги')
    r=years[0]
    metrics_cards=[('8,8 млн ₽','Выручка 2027','6000 га-проходов + 4000 обследований'),(f'{r["ebitda_rub"]/r["revenue_rub"]*100:.1f}%'.replace('.',','),'Маржа EBITDA 2027','До амортизации и налога УСН'),('1,432 млн ₽','Свободный CF 2027','После налога; до выплат инвесторам'),('13 000','Потолок га-проходов','100 дней × 130 га; подтвердить пилотом'),('10 клиентов','Портфель первого сезона','Договоры; календарь; актирование'),('0 работ','Без требуемой готовности','Учитывать также предотвращенные события')]
    for i,(num,title,body) in enumerate(metrics_cards):
        x=.55+(i%3)*4.15; y=1.88+(i//3)*2.25
        d.box(s,x,y,3.93,2.06,WHITE,True); d.text(s,num,x+.2,y+.18,3.53,.58,29,FOREST,True); d.text(s,title,x+.2,y+.83,3.53,.43,18,FOREST,True); d.text(s,body,x+.2,y+1.36,3.53,.61,16)
    d.text(s,'KPI подтверждаются актами, журналами, банковской выпиской и план-фактом — не презентацией.',.65,6.61,12,.43,16,MUTED)

    s=d.slide(22,'[ДОПУЩЕНИЕ] Риск-рейтинг экспертный; инвестиционное решение условное')
    d.table(s,['Риск / триггер','Ответ / владелец'],[['Нет правового основания','NO-GO; директор + юрист'],['Погода / конфликт окон','Перепланирование; опер. руководитель'],['Заказы ниже плана / цена падает','Пересчет NPV; директор'],['Поломка / ошибка внесения','Резерв, SOP, остановка; опер. руководитель'],['Перегрузка / дефицит ликвидности','Лимит заказов, кассовый план; директор']],x=.55,y=1.89,w=7.29,h=3.82,widths=[3.65,3.64],size=17)
    d.card(s,8.08,1.89,4.7,3.82,'INVESTMENT GATE','1  Письменная правовая карта\n2  Заявки на доступные 6 000 га\n3  Подтвержденные КП и кадры\n4  Кассовый план и резерв\n5  Приемка + апрельский пилот',size=18)
    d.banner(s,'Нет готовности → отложить капиталоемкий запуск. Мониторинг-only не доказан как прибыльная замена',6.05)
    d.text(s,'Масштабирование — после фактического сезона, проверки претензий, загрузки и продлений.',.65,6.69,12,.36,15,MUTED)
    d.finish()

if __name__=='__main__':
    build()
